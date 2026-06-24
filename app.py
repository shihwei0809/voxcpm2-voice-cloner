#!/usr/bin/env python3
"""
app.py - VoxCPM2 語音錄製工具
唯一的 UI：錄音 + 儲存。克隆與生成交由 AI Agent 處理。

用法：
  python app.py              # http://127.0.0.1:7860
"""

import os
import numpy as np
import gradio as gr
import json
import re

REPO_DIR = os.path.dirname(os.path.abspath(__file__))
SAMPLE_RATE = 16000

CUSTOM_CSS = """
body {
    font-family: 'Segoe UI', -apple-system, BlinkMacSystemFont, 'Roboto', 'Microsoft JhengHei', sans-serif;
}

/* Light Mode Card Styling */
.step-box {
    background: #ffffff !important;
    color: #1e293b !important;
    border-radius: 16px !important;
    padding: 24px !important;
    margin: 16px 0 !important;
    border: 1px solid #e2e8f0 !important;
    box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.05), 0 2px 4px -1px rgba(0, 0, 0, 0.03) !important;
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
}
.step-box:hover {
    box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.08), 0 4px 6px -2px rgba(0, 0, 0, 0.04) !important;
}

/* Dark Mode Card Styling */
.dark .step-box {
    background: #1e293b !important;
    color: #f8fafc !important;
    border: 1px solid #334155 !important;
    box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.2), 0 2px 4px -1px rgba(0, 0, 0, 0.1) !important;
}
.dark .step-box:hover {
    box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.3), 0 4px 6px -2px rgba(0, 0, 0, 0.2) !important;
}

/* Card titles and headings color override */
.step-box h2, .step-box p, .step-box li {
    color: inherit !important;
}
.step-box h2 {
    font-weight: 700 !important;
    margin-bottom: 12px !important;
}

/* Sample Text Box (yellow paper look) */
.sample-text-box {
    background: #fffbeb !important;
    color: #78350f !important;
    border: 1px solid #fef3c7 !important;
    padding: 20px !important;
    border-radius: 12px !important;
    font-size: 1.15em !important;
    line-height: 1.8 !important;
    margin: 12px 0 !important;
    font-weight: 500 !important;
    box-shadow: inset 0 2px 4px 0 rgba(0, 0, 0, 0.02) !important;
}
.dark .sample-text-box {
    background: #1e1b4b !important;
    color: #c7d2fe !important;
    border: 1px solid #312e81 !important;
}

/* Hint and helper text color */
.hint-text {
    text-align: center !important;
    margin-top: 12px !important;
    font-size: 0.95em !important;
    color: #475569 !important;
    font-weight: 500 !important;
}
.dark .hint-text {
    color: #94a3b8 !important;
}

/* Main Title styling */
.main-title {
    font-size: 2.2em !important;
    margin-bottom: 8px !important;
    font-weight: 800 !important;
    color: #0f172a !important;
}
.dark .main-title {
    color: #f8fafc !important;
}
.main-subtitle {
    font-size: 1.1em !important;
    color: #475569 !important;
}
.dark .main-subtitle {
    color: #94a3b8 !important;
}

/* Footer styling */
.footer-text {
    text-align: center !important;
    margin-top: 32px !important;
    padding: 16px !important;
    font-size: 0.85em !important;
    color: #94a3b8 !important;
    border-top: 1px solid #e2e8f0 !important;
}
.dark .footer-text {
    color: #64748b !important;
    border-top: 1px solid #334155 !important;
}
.footer-text a {
    color: #3b82f6 !important;
    text-decoration: none !important;
}
.footer-text a:hover {
    text-decoration: underline !important;
}
.success-msg { color:#10b981; font-weight:bold; }
"""


def get_gdrive_root():
    possible_roots = [
        r"G:\我的雲端硬碟\GOOGLE ANGET",
        r"G:\My Drive\GOOGLE ANGET"
    ]
    for root in possible_roots:
        if os.path.exists(root):
            return root
    return None


def get_gdrive_output_dir():
    root = get_gdrive_root()
    if root:
        gdrive_output = os.path.join(root, "AI 克隆聲音", "output")
        if os.path.exists(gdrive_output):
            return gdrive_output
    return None


def sync_voice_to_gdrive(voice_name, local_ref_path, local_prompt_path):
    root = get_gdrive_root()
    if root:
        try:
            gdrive_voice_dir = os.path.join(root, "AI 克隆聲音", "voices", voice_name)
            os.makedirs(gdrive_voice_dir, exist_ok=True)
            import shutil
            shutil.copy2(local_ref_path, os.path.join(gdrive_voice_dir, "ref_voice.wav"))
            shutil.copy2(local_prompt_path, os.path.join(gdrive_voice_dir, "prompt.txt"))
            print(f"已同步參考聲音「{voice_name}」至雲端硬碟庫")
            return True
        except Exception as e:
            print(f"同步參考聲音至雲端硬碟失敗: {e}")
    return False


def list_voices():
    local_vdir = os.path.join(REPO_DIR, "voices")
    voice_names = set()
    if os.path.exists(local_vdir):
        for d in os.listdir(local_vdir):
            if os.path.isdir(os.path.join(local_vdir, d)) and os.path.exists(os.path.join(local_vdir, d, "ref_voice.wav")):
                voice_names.add(d)
                
    # 掃描雲端硬碟的聲音庫
    root = get_gdrive_root()
    if root:
        gdrive_vdir = os.path.join(root, "AI 克隆聲音", "voices")
        if os.path.exists(gdrive_vdir):
            for d in os.listdir(gdrive_vdir):
                if os.path.isdir(os.path.join(gdrive_vdir, d)) and os.path.exists(os.path.join(gdrive_vdir, d, "ref_voice.wav")):
                    voice_names.add(d)
                    
    return sorted(list(voice_names))


def load_voice_details(voice_name):
    if not voice_name:
        return None, ""
        
    # 優先查本機
    wav_path = os.path.join(REPO_DIR, "voices", voice_name, "ref_voice.wav")
    prompt_path = os.path.join(REPO_DIR, "voices", voice_name, "prompt.txt")
    
    # 本機沒有，嘗試查雲端硬碟
    if not os.path.exists(wav_path):
        root = get_gdrive_root()
        if root:
            g_wav = os.path.join(root, "AI 克隆聲音", "voices", voice_name, "ref_voice.wav")
            g_txt = os.path.join(root, "AI 克隆聲音", "voices", voice_name, "prompt.txt")
            if os.path.exists(g_wav):
                wav_path = g_wav
                prompt_path = g_txt
                
    audio_val = wav_path if os.path.exists(wav_path) else None
    
    prompt_val = ""
    if os.path.exists(prompt_path):
        try:
            with open(prompt_path, "r", encoding="utf-8") as f:
                prompt_val = f.read().strip()
        except Exception:
            pass
            
    return audio_val, prompt_val


def load_sample_text():
    path = os.path.join(REPO_DIR, "texts", "sample_text.txt")
    if os.path.exists(path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read().strip()
    return ""


def list_generated_files():
    out_dir = os.path.join(REPO_DIR, "output")
    if not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)
        
    local_files = [f for f in os.listdir(out_dir) if f.endswith(".wav")]
    all_files = {f: os.path.join(out_dir, f) for f in local_files}
    
    # 掃描雲端硬碟已生成目錄
    gdrive_dir = get_gdrive_output_dir()
    if gdrive_dir:
        gdrive_files = [f for f in os.listdir(gdrive_dir) if f.endswith(".wav")]
        for f in gdrive_files:
            g_path = os.path.join(gdrive_dir, f)
            if f not in all_files:
                all_files[f] = g_path
            else:
                # 若檔名相同，選擇修改時間較新的
                if os.path.getmtime(g_path) > os.path.getmtime(all_files[f]):
                    all_files[f] = g_path
                    
    sorted_files = sorted(
        all_files.keys(),
        key=lambda x: os.path.getmtime(all_files[x]),
        reverse=True
    )
    return sorted_files


def load_history_audio(filename):
    if not filename:
        return None
    # 本機優先
    local_path = os.path.join(REPO_DIR, "output", filename)
    if os.path.exists(local_path):
        return local_path
    # 雲端硬碟
    gdrive_dir = get_gdrive_output_dir()
    if gdrive_dir:
        g_path = os.path.join(gdrive_dir, filename)
        if os.path.exists(g_path):
            return g_path
    return None



HTML_TEMPLATE = """<!DOCTYPE html>
<html lang="zh-TW">
<head>
  <meta charset="UTF-8">
  <title>{title}</title>
  <style>
    :root {
      --bg: #070d1a;
      --card: #0f1c2e;
      --surface: #162338;
      --accent: #38bdf8;
      --text: #f0f6ff;
      --muted: #8899aa;
      --border: rgba(56,189,248,0.15);
    }
    body {
      background: var(--bg);
      color: var(--text);
      font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
      height: 100vh;
      display: flex;
      flex-direction: column;
      align-items: center;
      justify-content: center;
      margin: 0;
      overflow: hidden;
    }
    .deck {
      width: 90vw;
      max-width: 960px;
      aspect-ratio: 16/9;
      background: var(--card);
      border: 1px solid var(--border);
      border-radius: 16px;
      display: flex;
      flex-direction: column;
      overflow: hidden;
      position: relative;
      box-shadow: 0 20px 50px rgba(0,0,0,0.5);
    }
    .slides-area {
      flex: 1;
      position: relative;
      background: #000;
    }
    .slide-img {
      width: 100%;
      height: 100%;
      object-fit: contain;
      position: absolute;
      opacity: 0;
      transition: opacity 0.5s ease;
    }
    .slide-img.active {
      opacity: 1;
    }
    .controls {
      height: 60px;
      background: rgba(17, 24, 39, 0.95);
      border-top: 1px solid var(--border);
      display: flex;
      align-items: center;
      justify-content: space-between;
      padding: 0 20px;
      gap: 16px;
    }
    .btn {
      background: rgba(255,255,255,0.05);
      border: 1px solid var(--border);
      color: var(--text);
      padding: 8px 16px;
      border-radius: 6px;
      cursor: pointer;
      font-weight: 600;
      display: flex;
      align-items: center;
      gap: 6px;
      transition: all 0.2s;
    }
    .btn:hover {
      background: var(--accent);
      color: #000;
      box-shadow: 0 0 10px rgba(56,189,248,0.4);
    }
    .progress-bar {
      flex: 1;
      height: 8px;
      background: var(--surface);
      border-radius: 4px;
      position: relative;
      cursor: pointer;
    }
    .progress-fill {
      height: 100%;
      background: linear-gradient(90deg, var(--accent), #818cf8);
      width: 0;
      border-radius: 4px;
      box-shadow: 0 0 8px var(--accent);
    }
    .time-label {
      font-size: 12px;
      color: var(--muted);
      font-family: monospace;
    }
  </style>
</head>
<body>
  <div class="deck">
    <div class="slides-area" id="slides-area">
      {slides_html}
    </div>
    <div class="controls">
      <button class="btn" id="prev-btn">◀ 上一頁</button>
      <button class="btn" id="play-btn">▶ 播放</button>
      <button class="btn" id="next-btn">下一頁 ▶</button>
      <div class="progress-bar" id="progress-bar">
        <div class="progress-fill" id="progress-fill"></div>
      </div>
      <span class="time-label" id="time-label">0:00 / 0:00</span>
    </div>
  </div>

  <audio id="audio" src="{audio_src}"></audio>

  <script>
    const timings = {timings_json};
    const slideImgs = document.querySelectorAll(".slide-img");
    const audio = document.getElementById("audio");
    const playBtn = document.getElementById("play-btn");
    const prevBtn = document.getElementById("prev-btn");
    const nextBtn = document.getElementById("next-btn");
    const progressFill = document.getElementById("progress-fill");
    const progressBar = document.getElementById("progress-bar");
    const timeLabel = document.getElementById("time-label");

    let currentSlide = 0;
    const startTimes = [0];
    for(let i=0; i<timings.length; i++) {{
      startTimes.push(startTimes[i] + timings[i]);
    }}
    const totalDuration = startTimes[startTimes.length - 1];

    function formatTime(s) {{
      const m = Math.floor(s / 60);
      const sec = Math.floor(s % 60).toString().padStart(2, '0');
      return `${m}:${sec}`;
    }}

    function updateUI() {{
      const t = audio.currentTime;
      let active = 0;
      for(let i=0; i<startTimes.length - 1; i++) {{
        if(t >= startTimes[i] && t < startTimes[i+1]) {{
          active = i;
          break;
        }}
      }}
      if (t >= totalDuration - 0.1) {{
        active = startTimes.length - 2;
      }}
      
      if(active !== currentSlide) {{
        slideImgs[currentSlide].classList.remove("active");
        slideImgs[active].classList.add("active");
        currentSlide = active;
      }}

      const pct = (t / totalDuration) * 100;
      progressFill.style.width = pct + "%";
      timeLabel.innerText = `${formatTime(t)} / ${formatTime(totalDuration)}`;
    }}

    audio.addEventListener("timeupdate", updateUI);
    
    playBtn.addEventListener("click", () => {{
      if(audio.paused) {{
        audio.play();
        playBtn.innerText = "⏸ 暫停";
      }} else {{
        audio.pause();
        playBtn.innerText = "▶ 播放";
      }}
    }});

    audio.addEventListener("ended", () => {{
      playBtn.innerText = "▶ 播放";
    }});

    prevBtn.addEventListener("click", () => {{
      if(currentSlide > 0) {{
        audio.currentTime = startTimes[currentSlide - 1];
        updateUI();
      }}
    }});

    nextBtn.addEventListener("click", () => {{
      if(currentSlide < startTimes.length - 2) {{
        audio.currentTime = startTimes[currentSlide + 1];
        updateUI();
      }}
    }});

    progressBar.addEventListener("click", (e) => {{
      const rect = progressBar.getBoundingClientRect();
      const pct = (e.clientX - rect.left) / rect.width;
      audio.currentTime = pct * totalDuration;
      updateUI();
    }});

    if (slideImgs.length > 0) {
      slideImgs[0].classList.add("active");
      const imgObj = new Image();
      imgObj.src = slideImgs[0].src;
      imgObj.onload = function() {
        const ratio = imgObj.naturalWidth / imgObj.naturalHeight;
        const deck = document.querySelector(".deck");
        if (deck && ratio) {
          deck.style.aspectRatio = ratio;
        }
      };
    }
  </script>
</body>
</html>
"""


def sync_html_and_slides_to_gdrive(html_path, slides_dir):
    root = get_gdrive_root()
    if root:
        try:
            import shutil
            # 複製 HTML
            gdrive_html = os.path.join(root, "AI 克隆聲音", os.path.basename(html_path))
            shutil.copy2(html_path, gdrive_html)
            
            # 複製投影片圖片目錄
            gdrive_slides_dir = os.path.join(root, "AI 克隆聲音", "output", os.path.basename(slides_dir))
            if os.path.exists(gdrive_slides_dir):
                shutil.rmtree(gdrive_slides_dir)
            shutil.copytree(slides_dir, gdrive_slides_dir)
            print("已同步 HTML 與投影片目錄至雲端硬碟")
            return True
        except Exception as e:
            print(f"同步投影片檔案至雲端硬碟失敗: {e}")
    return False

def extract_presentation_script(pptx_file):
    """Step 1: Extract narration text from PPT/PPTX/PDF, return as editable string."""
    if not pptx_file:
        return "⚠️ 請先上傳簡報/PDF 檔案。"
    try:
        base_name, ext = os.path.splitext(os.path.basename(pptx_file))
        ext = ext.lower()
        slide_texts = []

        if ext == ".pdf":
            import fitz
            doc = fitz.open(pptx_file)
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if not text:
                    text = f"這是第 {i} 頁 PDF 內容。"
                slide_texts.append(text)
        elif ext in [".pptx", ".ppt"]:
            if ext == ".pptx":
                try:
                    from pptx import Presentation as PptxPres
                    prs = PptxPres(pptx_file)
                    for i, slide in enumerate(prs.slides, 1):
                        note = ""
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                            note = slide.notes_slide.notes_text_frame.text.strip()
                        if not note:
                            # fallback: collect visible text from shapes
                            parts = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    t = shape.text_frame.text.strip()
                                    if t:
                                        parts.append(t)
                            note = " ".join(parts) if parts else f"這是第 {i} 頁投影片內容。"
                        slide_texts.append(note)
                except Exception as e:
                    return f"❌ 讀取 PPTX 備忘錄失敗：{e}"
            else:
                return "⚠️ .ppt 格式需要已安裝 PowerPoint，建議另存為 .pptx 後再上傳。"
        else:
            return "⚠️ 不支援的檔案格式。只支援 .pptx, .ppt 與 .pdf。"

        if not slide_texts:
            return "⚠️ 無法解析投影片或 PDF 頁數。"

        # Format as editable block: === 第 N 頁 === \n <text>
        lines = []
        for i, text in enumerate(slide_texts, 1):
            lines.append(f"=== 第 {i} 頁 ===")
            lines.append(text.strip())
            lines.append("")  # blank line separator
        return "\n".join(lines)
    except Exception as e:
        return f"❌ 擷取腳本失敗：{e}"


def script_text_to_slides(script_text):
    """Parse the editable script back into per-slide list."""
    slides = []
    current = []
    for line in script_text.splitlines():
        if line.strip().startswith("===") and line.strip().endswith("==="):
            if current:
                slides.append("\n".join(current).strip())
                current = []
        else:
            current.append(line)
    if current:
        slides.append("\n".join(current).strip())
    # Remove empty entries
    slides = [s for s in slides if s]
    return slides


def auto_synthesize_presentation(pptx_file, voice_name, script_text, progress=gr.Progress()):
    if not pptx_file:
        return "⚠️ 請先上傳簡報/PDF 檔案。"
    if not voice_name:
        return "⚠️ 請先選擇克隆配音聲音。"
    if not script_text or not script_text.strip():
        return "⚠️ 請先點擊「預覽腳本」，確認腳本內容後再生成。"

    try:

        base_name, ext = os.path.splitext(os.path.basename(pptx_file))
        ext = ext.lower()
        
        # 清理特殊字元
        import re
        base_name = re.sub(r'[^\w\s一-鿿-]', '', base_name).strip().replace(" ", "_")
        if not base_name:
            base_name = "presentation"
            
        slides_tmp_dir = os.path.join(REPO_DIR, "output", f"slides_{base_name}")
        os.makedirs(slides_tmp_dir, exist_ok=True)
        
        # --- 使用使用者確認/修改後的腳本，不重新擷取文字 ---
        slide_texts = script_text_to_slides(script_text)
        png_files = []

        # 仍需從原始檔匯出 PNG 投影片圖片
        if ext == ".pdf":
            progress(0, desc="正在讀取 PDF 檔案並匯出投影片圖片...")
            import fitz
            doc = fitz.open(pptx_file)
            total_slides_file = len(doc)
            for i, page in enumerate(doc, 1):
                progress((i / total_slides_file) * 0.15, desc=f"匯出第 {i}/{total_slides_file} 頁 PDF 為 PNG...")
                zoom = 2
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                png_name = f"slide_{i:03d}.png"
                out_png = os.path.join(slides_tmp_dir, png_name)
                pix.save(out_png)
                png_files.append(f"output/slides_{base_name}/{png_name}")

        elif ext in [".pptx", ".ppt"]:
            progress(0, desc="正在匯出投影片圖片...")
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            try:
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
                ppt_app.Visible = True
                prs_com = ppt_app.Presentations.Open(os.path.abspath(pptx_file), ReadOnly=True, Untitled=False, WithWindow=False)
                total_slides_file = prs_com.Slides.Count
                for i in range(1, total_slides_file + 1):
                    progress((i / total_slides_file) * 0.15, desc=f"匯出第 {i}/{total_slides_file} 頁投影片為 PNG...")
                    png_name = f"slide_{i:03d}.png"
                    out_png = os.path.join(slides_tmp_dir, png_name)
                    prs_com.Slides(i).Export(out_png, "PNG", 1920, 1080)
                    png_files.append(f"output/slides_{base_name}/{png_name}")
                prs_com.Close()
                ppt_app.Quit()
            finally:
                pythoncom.CoUninitialize()
        else:
            return "⚠️ 不支援的檔案格式。只支援 .pptx, .ppt 與 .pdf。"

        total_slides = len(slide_texts)
        if total_slides == 0:
            return "⚠️ 腳本內容為空，請確認腳本格式正確（每頁以 === 第 N 頁 === 分隔）。"
            
        progress(0.15, desc=f"共 {total_slides} 頁。正在載入 VoxCPM2 語音模型...")
        
        from voxcpm import VoxCPM
        import soundfile as sf
        import numpy as np
        
        # 決定裝置
        import torch
        if torch.cuda.is_available():
            dev = "cuda"
        elif hasattr(torch, "xpu") and torch.xpu.is_available():
            dev = "xpu"
        else:
            dev = "cpu"
            
        ref_wav, prompt_text = load_voice_details(voice_name)
        if not ref_wav:
            return f"⚠️ 找不到克隆聲音「{voice_name}」的參考音檔。"
            
        model = VoxCPM.from_pretrained('openbmb/VoxCPM2', load_denoiser=False, device=dev, optimize=False)
        
        # 合成每張投影片語音
        wav_clips = []
        timings = []
        
        for i, note in enumerate(slide_texts, 1):
            pct = 0.2 + (i / total_slides) * 0.55
            progress(pct, desc=f"[{i}/{total_slides}] 正在合成第 {i} 頁備忘錄語音... ({len(note)} 字)")
            wav = model.generate(
                text=note,
                prompt_wav_path=ref_wav,
                prompt_text=prompt_text,
                reference_wav_path=ref_wav,
                cfg_value=1.5,
                inference_timesteps=10,
            )
            wav_clips.append(wav)
            dur = len(wav) / model.tts_model.sample_rate
            timings.append(dur)
            
            # 每頁結束加停頓 0.5 秒
            pause = np.zeros(int(0.5 * model.tts_model.sample_rate), dtype=wav.dtype)
            wav_clips.append(pause)
            timings[-1] += 0.5
            
        progress(0.8, desc="正在合併並儲存完整配音檔...")
        full_audio = np.concatenate(wav_clips)
        
        out_wav_name = f"{base_name}_{voice_name}.wav"
        out_wav_path = os.path.join(REPO_DIR, "output", out_wav_name)
        os.makedirs(os.path.dirname(out_wav_path), exist_ok=True)
        sf.write(out_wav_path, full_audio, model.tts_model.sample_rate)
        
        # 雲端同步
        sync_wav_dest = ""
        root = get_gdrive_root()
        if root:
            try:
                gdrive_output_dir = os.path.join(root, "AI 克隆聲音", "output")
                os.makedirs(gdrive_output_dir, exist_ok=True)
                import shutil
                dest_path = os.path.join(gdrive_output_dir, out_wav_name)
                shutil.copy2(out_wav_path, dest_path)
                sync_wav_dest = dest_path
            except Exception as e:
                print(f"同步 wav 失敗: {e}")
 
        progress(0.85, desc="正在產生互動式簡報網頁...")
        slides_html_parts = []
        for i, img_path in enumerate(png_files):
            slides_html_parts.append(f'<img class="slide-img" src="{img_path}">')
        slides_html = "\n      ".join(slides_html_parts)
        
        # 內嵌 wav 音軌 base64
        import base64
        with open(out_wav_path, "rb") as af:
            wav_b64 = base64.b64encode(af.read()).decode("ascii")
        audio_uri = f"data:audio/wav;base64,{wav_b64}"
        
        html_content = HTML_TEMPLATE.replace("{title}", base_name)
        html_content = html_content.replace("{slides_html}", slides_html)
        html_content = html_content.replace("{audio_src}", audio_uri)
        html_content = html_content.replace("{timings_json}", json.dumps(timings))
        
        out_html_name = f"{base_name}.html"
        out_html_path = os.path.join(REPO_DIR, out_html_name)
        with open(out_html_path, "w", encoding="utf-8") as f:
            f.write(html_content)
            
        progress(0.92, desc="同步網頁與投影片檔案至雲端硬碟...")
        sync_html_and_slides_to_gdrive(out_html_path, slides_tmp_dir)
        
        progress(1.0, desc="完成！")
        
        return (
            f"🎉 投影片自動合成成功！\n\n"
            f"1. 完整配音已存至本機: `output/{out_wav_name}`\n"
            f"2. 互動式網頁已產生: `{out_html_name}`\n"
            f"3. 投影片底圖資料夾: `output/slides_{base_name}/`\n\n"
            f"🔗 雲端硬碟同步：\n"
            f"   - 音檔: `{sync_wav_dest or '無（未偵測到雲端硬碟）'}`\n"
            f"   - 網頁: `G:\\我的雲端硬碟\\GOOGLE ANGET\\AI 克隆聲音\\{out_html_name}`\n\n"
            f"👉 請切換到「📊 互動式簡報播放」分頁，點擊「🔄 重新掃描」，並在下拉選單選擇 `{out_html_name}` 即可開始播放對齊克隆語音的投影片！"
        )
    except Exception as e:
        import traceback
        err_msg = traceback.format_exc()
        return f"❌ 自動合成失敗：\n{e}\n\n詳細日誌：\n{err_msg}"


def scan_presentations():
    """掃描專案目錄下所有可用的簡報 HTML 檔案（排除 _embedded 結尾）。"""
    import glob
    pattern = os.path.join(REPO_DIR, "*.html")
    files = sorted(glob.glob(pattern))
    # 過濾掉 _embedded 版本（內部使用），只顯示原始簡報
    names = [
        os.path.basename(f) for f in files
        if not os.path.basename(f).endswith("_embedded.html")
    ]
    return names if names else ["（尚無簡報）"]


def load_presentation_iframe(filename: str) -> str:
    """載入指定簡報 HTML，自動嵌入同目錄下同名 WAV（若存在），回傳 iframe HTML。"""
    import base64, html as html_mod, glob

    if not filename or filename == "（尚無簡報）":
        return "<div style='padding:40px; text-align:center; color:#64748b;'>請先選擇一份簡報。</div>"

    base = os.path.splitext(filename)[0]          # e.g. summary_presentation
    html_path     = os.path.join(REPO_DIR, filename)
    embedded_path = os.path.join(REPO_DIR, base + "_embedded.html")

    # 若已有 embedded 版直接用
    if os.path.exists(embedded_path):
        path = embedded_path
    elif os.path.exists(html_path):
        # 嘗試找同名 WAV（掃描本機與雲端硬碟的 output/ 目錄）
        wav_candidates = []
        local_dir = os.path.join(REPO_DIR, "output")
        if os.path.exists(local_dir):
            wav_candidates.extend(glob.glob(os.path.join(local_dir, "*.wav")))
            
        gdrive_dir = get_gdrive_output_dir()
        if gdrive_dir:
            wav_candidates.extend(glob.glob(os.path.join(gdrive_dir, "*.wav")))
            
        # 移除重複的檔名，選擇時間較新的
        unique_candidates = {}
        for w in wav_candidates:
            fname = os.path.basename(w)
            if fname not in unique_candidates:
                unique_candidates[fname] = w
            else:
                if os.path.getmtime(w) > os.path.getmtime(unique_candidates[fname]):
                    unique_candidates[fname] = w
        wav_candidates = list(unique_candidates.values())
        
        # 優先找名稱最接近的
        matched_wav = None
        for w in wav_candidates:
            if base.lower().replace("_", "") in os.path.basename(w).lower().replace("_", ""):
                matched_wav = w
                break
        if not matched_wav and wav_candidates:
            # 依時間遞減排序，拿最新的一個
            wav_candidates.sort(key=os.path.getmtime, reverse=True)
            matched_wav = wav_candidates[0]   # fallback: 第一個 wav

        with open(html_path, "r", encoding="utf-8") as f:
            html_content = f.read()

        if matched_wav:
            with open(matched_wav, "rb") as af:
                b64 = base64.b64encode(af.read()).decode("ascii")
            data_uri = f"data:audio/wav;base64,{b64}"
            # 替換所有 .wav 路徑為 base64
            import re
            html_content = re.sub(r'src="[^"]*\.wav"', f'src="{data_uri}"', html_content)

        with open(embedded_path, "w", encoding="utf-8") as f:
            f.write(html_content)
        path = embedded_path
    else:
        return f"<div style='padding:20px; color:#ef4444;'>⚠️ 找不到檔案：{filename}</div>"

    with open(path, "r", encoding="utf-8") as f:
        html_content = f.read()
    escaped = html_mod.escape(html_content)
    return f'<iframe srcdoc="{escaped}" style="width:100%; height:700px; border:none; border-radius:16px;"></iframe>'



def on_save(audio_mic, audio_upload, voice_name):
    audio_data = audio_mic or audio_upload
    if audio_data is None:
        return "⚠️ 請先按「開始錄音」錄下你的聲音。"
    if not voice_name or not voice_name.strip():
        return "⚠️ 請為你的聲音取一個名字。"

    name = voice_name.strip()
    try:
        sr, audio = audio_data
        if audio.ndim > 1:
            audio = audio.mean(axis=1)
        if sr != SAMPLE_RATE:
            import resampy
            audio = resampy.resample(audio, sr, SAMPLE_RATE)
        peak = np.abs(audio).max()
        if peak > 0 and peak < 0.01:
            return "⚠️ 錄到的音量太小，請靠近麥克風再錄一次。"
        if peak > 0:
            audio = audio / peak * 0.95

        vdir = os.path.join(REPO_DIR, "voices", name)
        os.makedirs(vdir, exist_ok=True)
        import soundfile as sf
        ref_path = os.path.join(vdir, "ref_voice.wav")
        prompt_path = os.path.join(vdir, "prompt.txt")
        sf.write(ref_path, audio.astype(np.float32), SAMPLE_RATE)
        with open(prompt_path, "w", encoding="utf-8") as f:
            f.write(load_sample_text())

        # 同步錄製的聲音到雲端硬碟共用資料庫
        sync_voice_to_gdrive(name, ref_path, prompt_path)

        duration = len(audio) / SAMPLE_RATE
        existing = list_voices()
        voice_list = "、".join(existing)
        return (
            f"✅ 聲音「{name}」錄製成功！（{duration:.0f} 秒）\n\n"
            f"目前已錄製的聲音：{voice_list}\n\n"
            f"👉 接下來，你可以對 AI 說：\n"
            f"「 用 {name} 的聲音說一段話 」"
        )
    except Exception as e:
        return f"❌ 儲存失敗：{e}"


def build_ui():
    sample_text = load_sample_text()

    with gr.Blocks(title="VoxCPM2 語音與簡報工具", css=CUSTOM_CSS) as app:
        gr.HTML("""
        <div style="text-align:center; margin-bottom:24px;">
          <h1 class="main-title">🎙️ VoxCPM2 語音與簡報工具</h1>
          <p class="main-subtitle">在此錄製您的參考聲音、管理音庫，並以互動式網報同步播放您的專案簡報與克隆配音。</p>
        </div>
        """)

        with gr.Tabs():
            with gr.Tab("🎙️ 聲音錄製與庫管理"):
                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## ✏️ 為聲音取名字")
                    voice_name_input = gr.Textbox(
                        label="",
                        placeholder="例如：王老師、林主任...",
                        show_label=False,
                    )

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 📖 請念這段文字")
                    gr.HTML(
                        f'<div class="sample-text-box">{sample_text}</div>'
                    )

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 🎤 錄音並儲存")
                    audio_mic = gr.Audio(label="", type="numpy", sources=["microphone"], show_label=False)

                    gr.HTML('<div class="hint-text">👆 錄完後按這裡儲存 👇</div>')

                    with gr.Row():
                        save_btn = gr.Button("⬇️ 儲存聲音", variant="primary", size="lg")

                    with gr.Accordion("或上傳已錄好的音檔", open=False):
                        audio_upload = gr.Audio(label="", type="numpy", sources=["upload"], show_label=False)

                save_msg = gr.Textbox(label="", show_label=False, lines=5, interactive=False)

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 🎙️ 已錄製的參考聲音 (本機聲音庫)")
                    
                    voice_choices = list_voices()
                    voice_val = voice_choices[0] if voice_choices else None
                    voice_audio, voice_prompt = load_voice_details(voice_val)
                    
                    with gr.Row():
                        voice_dropdown = gr.Dropdown(
                            label="選擇聲音名稱",
                            choices=voice_choices,
                            value=voice_val,
                            interactive=True,
                            scale=4,
                            show_label=False
                        )
                        refresh_voices_btn = gr.Button("🔄 重新整理", variant="secondary", scale=1)
                        
                    with gr.Row():
                        voice_audio_player = gr.Audio(
                            label="播放參考音",
                            value=voice_audio,
                            type="filepath",
                            interactive=False,
                            scale=2
                        )
                        voice_prompt_display = gr.Textbox(
                            label="朗讀文字 (Prompt)",
                            value=voice_prompt,
                            interactive=False,
                            scale=3,
                            lines=3
                        )

                def on_voice_change(name):
                    return load_voice_details(name)
                    
                voice_dropdown.change(
                    fn=on_voice_change,
                    inputs=[voice_dropdown],
                    outputs=[voice_audio_player, voice_prompt_display]
                )
                
                def on_refresh_voices():
                    choices = list_voices()
                    val = choices[0] if choices else None
                    audio, prompt = load_voice_details(val)
                    return gr.Dropdown(choices=choices, value=val), audio, prompt

                refresh_voices_btn.click(
                    fn=on_refresh_voices,
                    outputs=[voice_dropdown, voice_audio_player, voice_prompt_display]
                )

                save_btn.click(
                    fn=on_save,
                    inputs=[audio_mic, audio_upload, voice_name_input],
                    outputs=[save_msg],
                ).then(
                    fn=on_refresh_voices,
                    outputs=[voice_dropdown, voice_audio_player, voice_prompt_display]
                )

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 🎵 已生成的複製語音紀錄 (AI 生成庫)")
                    
                    initial_files = list_generated_files()
                    initial_val = initial_files[0] if initial_files else None
                    initial_audio = load_history_audio(initial_val)
                    
                    with gr.Row():
                        history_dropdown = gr.Dropdown(
                            label="選擇音檔",
                            choices=initial_files,
                            value=initial_val,
                            interactive=True,
                            scale=4,
                            show_label=False
                        )
                        refresh_btn = gr.Button("🔄 重新整理", variant="secondary", scale=1)
                        
                    history_audio = gr.Audio(
                        label="播放音檔",
                        value=initial_audio,
                        type="filepath",
                        interactive=False
                    )

                def on_dropdown_change(filename):
                    return load_history_audio(filename)
                    
                history_dropdown.change(
                    fn=on_dropdown_change,
                    inputs=[history_dropdown],
                    outputs=[history_audio]
                )
                
                def on_refresh():
                    files = list_generated_files()
                    val = files[0] if files else None
                    audio_path = load_history_audio(val)
                    return gr.Dropdown(choices=files, value=val), audio_path

                refresh_btn.click(
                    fn=on_refresh,
                    outputs=[history_dropdown, history_audio]
                )

            with gr.Tab("📊 互動式簡報播放"):
                with gr.Row():
                    pres_dropdown = gr.Dropdown(
                        label="選擇簡報",
                        choices=scan_presentations(),
                        value=(scan_presentations() or [None])[0],
                        scale=5,
                    )
                    pres_refresh = gr.Button("🔄 重新掃描", variant="secondary", scale=1)

                initial_pres = (scan_presentations() or [None])[0]
                pres_iframe = gr.HTML(
                    value=load_presentation_iframe(initial_pres) if initial_pres else ""
                )

                pres_dropdown.change(
                    fn=load_presentation_iframe,
                    inputs=[pres_dropdown],
                    outputs=[pres_iframe],
                )

                def refresh_pres_list():
                    files = scan_presentations()
                    first = files[0] if files else None
                    return gr.Dropdown(choices=files, value=first), load_presentation_iframe(first)

                pres_refresh.click(
                    fn=refresh_pres_list,
                    outputs=[pres_dropdown, pres_iframe],
                )

            with gr.Tab("🏗️ 投影片自動合成"):
                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 1. 上傳簡報或 PDF 檔案 (.pptx, .ppt, .pdf)")
                    gr.Markdown("💡 *系統將自動擷取 PPT/PPTX **備忘錄** 或 PDF **頁面文字** 作為配音腳本，您可在下方確認並修改後再生成。*")
                    pptx_upload = gr.File(label="上傳簡報/PDF 檔案", file_types=[".pptx", ".ppt", ".pdf"], type="filepath")
                    extract_btn = gr.Button("📄 Step 1：預覽腳本", variant="secondary", size="lg")

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 2. 確認並修改腳本")
                    gr.Markdown("✏️ *每頁以 `=== 第 N 頁 ===` 分隔，您可以直接在下方文字框修改每頁的朗讀內容。*")
                    script_preview = gr.Textbox(
                        label="腳本內容（可直接編輯）",
                        lines=20,
                        interactive=True,
                        placeholder="點擊上方「預覽腳本」後，腳本將會顯示在此處，您可以直接修改每頁文字後再按生成..."
                    )

                with gr.Column(elem_classes="step-box"):
                    gr.Markdown("## 3. 選擇聲音並生成")
                    voice_dropdown_synth = gr.Dropdown(
                        label="選擇克隆配音聲音",
                        choices=list_voices(),
                        value=list_voices()[0] if list_voices() else None,
                        interactive=True
                    )
                    synth_btn = gr.Button("⚡ Step 2：確認腳本，開始合成語音 + 簡報", variant="primary", size="lg")

                synth_result = gr.Textbox(label="合成日誌與結果", lines=12, interactive=False)

                extract_btn.click(
                    fn=extract_presentation_script,
                    inputs=[pptx_upload],
                    outputs=[script_preview]
                )

                synth_btn.click(
                    fn=auto_synthesize_presentation,
                    inputs=[pptx_upload, voice_dropdown_synth, script_preview],
                    outputs=[synth_result]
                )

        gr.HTML("""
        <div class="footer-text">
        VoxCPM2（Apache-2.0 可商用） |
        <a href="https://github.com/mathruffian-dot/voxcpm2-voice-cloner" target="_blank">GitHub</a>
        </div>
        """)

    return app


if __name__ == "__main__":
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("--port", "-p", type=int, default=7860)
    args = p.parse_args()
    build_ui().launch(server_name="0.0.0.0", server_port=args.port, theme=gr.themes.Soft(), allowed_paths=[REPO_DIR])
